#!/usr/bin/env python3
"""
PowerPoint to Markdown Converter
Extracts content from .pptx files and converts to Markdown format
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys
import os
import re
from PIL import Image
import io

def clean_text(text):
    """Clean and normalize text content"""
    if not text:
        return ""
    # Remove extra whitespace and normalize line breaks
    text = re.sub(r'\s+', ' ', text.strip())
    return text

def extract_images_from_slide(slide, slide_number, output_dir):
    """Extract images from a slide and save them"""
    images = []
    
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                # Get the image data
                image = shape.image
                image_bytes = image.blob
                
                # Determine file extension
                ext = 'png'  # default
                if image.content_type == 'image/jpeg':
                    ext = 'jpg'
                elif image.content_type == 'image/png':
                    ext = 'png'
                elif image.content_type == 'image/gif':
                    ext = 'gif'
                
                # Create filename
                filename = f"slide_{slide_number:02d}_image_{len(images)+1}.{ext}"
                filepath = os.path.join(output_dir, filename)
                
                # Save image
                with open(filepath, 'wb') as f:
                    f.write(image_bytes)
                
                images.append({
                    'filename': filename,
                    'filepath': filepath,
                    'content_type': image.content_type
                })
                
                print(f"  📷 Extracted image: {filename}")
                
            except Exception as e:
                print(f"  ⚠️ Could not extract image from slide {slide_number}: {e}")
    
    return images

def extract_slide_content(slide):
    """Extract all text content from a slide"""
    content = {
        'title': '',
        'subtitle': '',
        'content': [],
        'notes': '',
        'images': []
    }
    
    # Extract text from all shapes
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text = clean_text(shape.text)
            if text:
                # Try to identify title vs content based on position and formatting
                if shape.top < 1000000:  # Rough heuristic for title position
                    if not content['title']:
                        content['title'] = text
                    elif not content['subtitle'] and len(text) < 100:
                        content['subtitle'] = text
                    else:
                        content['content'].append(text)
                else:
                    content['content'].append(text)
    
    # Extract speaker notes if present
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
        content['notes'] = clean_text(slide.notes_slide.notes_text_frame.text)
    
    return content

def format_as_markdown(slide_content, slide_number, total_slides, images=None):
    """Convert slide content to Markdown format"""
    md_content = []
    
    # Add slide header
    md_content.append(f"## Slide {slide_number} of {total_slides}")
    md_content.append("")
    
    # Add title
    if slide_content['title']:
        md_content.append(f"# {slide_content['title']}")
        md_content.append("")
    
    # Add subtitle
    if slide_content['subtitle']:
        md_content.append(f"*{slide_content['subtitle']}*")
        md_content.append("")
    
    # Add images if any
    if images:
        md_content.append("### Images")
        for img in images:
            md_content.append(f"![Slide {slide_number} Image](images/{img['filename']})")
            md_content.append("")
    
    # Add content
    for item in slide_content['content']:
        # Check if content looks like a bullet point
        if item.startswith('•') or item.startswith('-') or item.startswith('*'):
            md_content.append(f"- {item[1:].strip()}")
        else:
            # Check if it might be a list item without bullet
            if len(item) < 200 and not item.endswith('.'):
                md_content.append(f"- {item}")
            else:
                md_content.append(item)
        md_content.append("")
    
    # Add speaker notes if present
    if slide_content['notes']:
        md_content.append("### Speaker Notes")
        md_content.append(slide_content['notes'])
        md_content.append("")
    
    # Add separator
    md_content.append("---")
    md_content.append("")
    
    return "\n".join(md_content)

def convert_pptx_to_markdown(pptx_file, output_file=None):
    """Main function to convert PowerPoint to Markdown"""
    try:
        # Load presentation
        prs = Presentation(pptx_file)
        
        # Prepare output
        if output_file is None:
            output_file = os.path.splitext(pptx_file)[0] + '.md'
        
        # Create images directory
        output_dir = os.path.dirname(output_file)
        images_dir = os.path.join(output_dir, 'images')
        os.makedirs(images_dir, exist_ok=True)
        
        markdown_content = []
        total_images = 0
        
        # Add document header
        filename = os.path.basename(pptx_file)
        markdown_content.append(f"# {os.path.splitext(filename)[0]}")
        markdown_content.append("")
        markdown_content.append(f"*Converted from PowerPoint: {filename}*")
        markdown_content.append("")
        markdown_content.append(f"**Total Slides:** {len(prs.slides)}")
        markdown_content.append("")
        markdown_content.append("---")
        markdown_content.append("")
        
        # Process each slide
        for i, slide in enumerate(prs.slides, 1):
            print(f"Processing slide {i} of {len(prs.slides)}...")
            
            # Extract text content
            slide_content = extract_slide_content(slide)
            
            # Extract images
            slide_images = extract_images_from_slide(slide, i, images_dir)
            total_images += len(slide_images)
            
            # Format as markdown
            slide_md = format_as_markdown(slide_content, i, len(prs.slides), slide_images)
            markdown_content.append(slide_md)
        
        # Write to file
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write('\n'.join(markdown_content))
        
        print(f"✅ Successfully converted {pptx_file} to {output_file}")
        print(f"📄 Processed {len(prs.slides)} slides")
        print(f"📷 Extracted {total_images} images to {images_dir}")
        
        return output_file
        
    except Exception as e:
        print(f"❌ Error converting PowerPoint: {e}")
        return None

def main():
    """Command line interface"""
    if len(sys.argv) < 2:
        print("Usage: python pptx_to_md.py <input.pptx> [output.md]")
        sys.exit(1)
    
    input_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else None
    
    if not os.path.exists(input_file):
        print(f"❌ Error: File '{input_file}' not found")
        sys.exit(1)
    
    if not input_file.lower().endswith('.pptx'):
        print("❌ Error: Input file must be a .pptx file")
        sys.exit(1)
    
    convert_pptx_to_markdown(input_file, output_file)

if __name__ == "__main__":
    main()
